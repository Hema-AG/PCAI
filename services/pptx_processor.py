import os
import tempfile
from typing import List, Tuple
from pptx import Presentation
from PIL import Image


class PPTXProcessor:
    def validate_pptx_file(self, file_path: str) -> bool:
        """Validate if the file is a valid PPTX file"""
        try:
            presentation = Presentation(file_path)
            return len(presentation.slides) > 0
        except Exception:
            return False

    def extract_text_from_pptx(self, file_path: str) -> List[str]:
        """Extract text from each slide of the PPTX file"""
        presentation = Presentation(file_path)
        slide_texts = []

        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            slide_texts.append("\n".join(slide_text))

        return slide_texts

    def extract_slide_images(self, file_path: str, output_dir: str) -> List[str]:
        """Extract each slide as an image and return list of image paths"""
        presentation = Presentation(file_path)
        image_paths = []

        for i, slide in enumerate(presentation.slides):
            # In a real implementation, you would convert slides to images
            # This is a simplified version - you might need a different approach
            # depending on your specific requirements

            # For now, we'll create a placeholder image
            img_path = os.path.join(output_dir, f"slide_{i+1}.png")
            self._create_placeholder_image(img_path, f"Slide {i+1}")
            image_paths.append(img_path)

        return image_paths

    def _create_placeholder_image(self, path: str, text: str):
        """Create a placeholder image for demonstration"""
        from PIL import Image, ImageDraw, ImageFont

        # Create a blank image
        img = Image.new('RGB', (1280, 720), color=(73, 109, 137))
        d = ImageDraw.Draw(img)

        # Add text to the image
        try:
            font = ImageFont.truetype("arial.ttf", 40)
        except:
            font = ImageFont.load_default()

        d.text((100, 100), text, fill=(255, 255, 255), font=font)
        img.save(path)
